import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Pt

# Abre el archivo XML y parsea su contenido
tree = ET.parse('slide.xml')
root = tree.getroot()

# Crea una presentación PowerPoint
ppt = Presentation()

# Itera sobre los elementos del archivo XML y crea diapositivas en la presentación
for diapo in root.iter('diapositiva'):
    slide = ppt.slides.add_slide(ppt.slide_layouts[1]) # se usa el layout por defecto
    titulo = diapo.find('titulo').text
    
    # Agrega el título de la diapositiva
    title_shape = slide.shapes.title
    title_shape.text = titulo
    title_frame = title_shape.text_frame
    for paragraph in title_frame.paragraphs[0:]:
        paragraph.font.size = Pt(40) # establece el tamaño de fuente de los párrafos adicionales
    
    # Busca elementos de lista y los agrega a la diapositiva
    elementos = diapo.findall('elementoLista')
    if elementos:
        cuerpo = ''
        for elem in elementos:
            cuerpo += f'\n{elem.text}'
        cuerpo = cuerpo[1:] # eliminar el primer salto de línea
        # Agrega el cuerpo de la diapositiva
        body_shape = slide.placeholders[1]
        body_shape.text = cuerpo
        body_frame = body_shape.text_frame
        for paragraph in body_frame.paragraphs[0:]:
            paragraph.font.size = Pt(18) # establece el tamaño de fuente de los párrafos adicionales

# Guarda la presentación PowerPoint
ppt.save('presentacion.pptx')